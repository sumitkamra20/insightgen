# Below function should be included in this file
# pdf conversion to image files
# headline generation function that takes folder with image files as input
# and returns a list of headlines
#
# Final function that inputs the list of headline in pptx file and outputs the pptx file

# Imports
import os
from pdf2image import convert_from_path, convert_from_bytes
import logging
from io import BytesIO
import shutil
from pathlib import Path
from pptx import Presentation
import base64
from pptx.enum.shapes import PP_PLACEHOLDER
from typing import List, Dict, Union, Optional, BinaryIO, Tuple


def extract_slide_metadata(
    input_folder: str = None,
    pptx_file_content: bytes = None,
    pptx_filename: str = None
) -> dict:
    """
    Extracts metadata from each slide in a PPTX file.

    Args:
        input_folder (str, optional): Path to the folder containing the PPTX file.
        pptx_file_content (bytes, optional): PPTX file content as bytes.
        pptx_filename (str, optional): Name of the PPTX file when provided as bytes.

    Returns:
        dict: Dictionary storing slide metadata including layout, content status, placeholder availability,
              and placeholders for observations.
    """
    presentation = None

    # Handle file from disk
    if input_folder:
        # Find the PPTX file in the input folder
        pptx_files = [f for f in os.listdir(input_folder) if f.endswith('.pptx')]

        if not pptx_files:
            raise FileNotFoundError("No PPTX file found in the input folder.")
        if len(pptx_files) > 1:
            raise ValueError("Multiple PPTX files found. Please keep only one.")

        pptx_path = os.path.join(input_folder, pptx_files[0])
        presentation = Presentation(pptx_path)

    # Handle file from memory
    elif pptx_file_content:
        pptx_stream = BytesIO(pptx_file_content)
        presentation = Presentation(pptx_stream)

    else:
        raise ValueError("Either input_folder or pptx_file_content must be provided.")

    slide_data = {}

    # Iterate through slides and extract metadata
    for slide_number, slide in enumerate(presentation.slides, start=1):
        layout_name = slide.slide_layout.name  # Extract layout name

        # Mark slide as non-content if its layout name starts with "HEADER" (case-insensitive)
        content_slide = not layout_name.upper().startswith("HEADER")

        # Check if a title placeholder exists
        has_placeholder = any(
            shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE
            for shape in slide.shapes
        )

        # Initialize empty fields for later functions to fill
        slide_data[slide_number] = {
            "layout": layout_name,
            "content_slide": content_slide,
            "has_placeholder": has_placeholder,
            "key_observations": "",
            "slide_headline": "",
            "speaker_notes": "",
            "filename": pptx_filename if pptx_filename else (pptx_files[0] if input_folder else "presentation.pptx")
        }

    return slide_data

def generate_slide_images_base64(
    input_folder: str = None,
    slide_data: dict = None,
    pdf_file_content: bytes = None,
    img_format: str = "JPEG",
    dpi: int = 200
) -> dict:
    """
    Converts PDF slides to images, encodes them in base64, and updates the slide_data dictionary.
    Excludes non-content slides (e.g., Header or Divider) from image processing.

    Args:
        input_folder (str, optional): Directory containing input PDF and PPTX files.
        slide_data (dict): Dictionary storing slide metadata.
        pdf_file_content (bytes, optional): PDF file content as bytes.
        img_format (str): Image format (default: JPEG).
        dpi (int): Resolution for image conversion.

    Returns:
        dict: Updated slide metadata dictionary with base64 images (only for content slides).
    """
    logging.info("Starting PDF to image conversion...")

    if not slide_data:
        raise ValueError("slide_data must be provided")

    images = None

    # Handle file from disk
    if input_folder:
        # Validate input directory
        if not os.path.exists(input_folder):
            raise ValueError(f"Input directory does not exist: {input_folder}")

        # Find PDF file in the input folder
        pdf_files = [f for f in os.listdir(input_folder) if f.endswith('.pdf')]

        if not pdf_files:
            logging.error("No PDF file found in the input folder.")
            return slide_data

        if len(pdf_files) > 1:
            logging.error("Multiple PDF files found. Please keep only one.")
            return slide_data

        pdf_path = os.path.join(input_folder, pdf_files[0])

        # Convert PDF to images (in-memory)
        images = convert_from_path(pdf_path, dpi=dpi)
        logging.info(f"PDF successfully converted to {len(images)} images from file.")

    # Handle file from memory
    elif pdf_file_content:
        # Convert PDF bytes to images (in-memory)
        images = convert_from_bytes(pdf_file_content, dpi=dpi)
        logging.info(f"PDF successfully converted to {len(images)} images from bytes.")

    else:
        raise ValueError("Either input_folder or pdf_file_content must be provided.")

    # Process only content slides
    for i, image in enumerate(images, start=1):
        slide_number = i  # Assuming slides and PDF pages match 1:1

        # Skip non-content slides
        if slide_number not in slide_data or not slide_data[slide_number]["content_slide"]:
            slide_data[slide_number]["status"] = "Skipped (Non-content slide)"
            continue

        # Convert image to base64 (in-memory)
        img_byte_arr = BytesIO()
        image.save(img_byte_arr, format=img_format)
        img_byte_arr.seek(0)
        base64_image = base64.b64encode(img_byte_arr.read()).decode('utf-8')

        # Store base64 image in slide_data dictionary
        slide_data[slide_number]["image_base64"] = base64_image
        slide_data[slide_number]["status"] = "Image processed"

        logging.info(f"Slide {slide_number}: Image converted and stored as base64.")

    logging.info("Base64 images stored successfully in slide metadata.")

    return slide_data

def insert_headlines_into_pptx(
    input_folder: str = None,
    output_folder: str = None,
    slide_data: dict = None,
    pptx_file_content: bytes = None,
    save_as_new: bool = True
) -> Union[str, Tuple[str, bytes]]:
    """
    Inserts AI-generated headlines into slide title placeholders and observations into speaker notes.

    Args:
        input_folder (str, optional): Path to the folder containing the input PPTX file.
        output_folder (str, optional): Path to the folder where the modified PPTX file should be saved.
        slide_data (dict): Dictionary storing slide metadata, headlines, and observations.
        pptx_file_content (bytes, optional): PPTX file content as bytes.
        save_as_new (bool): Whether to save as a new file.

    Returns:
        Union[str, Tuple[str, bytes]]:
            - If input_folder and output_folder are provided: Path to the saved PowerPoint file.
            - If pptx_file_content is provided: Tuple of (filename, bytes) of the modified presentation.
    """
    logging.info("Starting headline and observations insertion into PowerPoint...")

    if not slide_data:
        raise ValueError("slide_data must be provided")

    presentation = None
    original_filename = None

    # Handle file from disk
    if input_folder and output_folder:
        pptx_files = [f for f in os.listdir(input_folder) if f.endswith('.pptx')]
        if not pptx_files:
            raise FileNotFoundError("No PPTX file found in the input folder.")
        if len(pptx_files) > 1:
            raise ValueError("Multiple PPTX files found. Please keep only one.")

        pptx_path = os.path.join(input_folder, pptx_files[0])
        presentation = Presentation(pptx_path)
        original_filename = pptx_files[0]

    # Handle file from memory
    elif pptx_file_content:
        pptx_stream = BytesIO(pptx_file_content)
        presentation = Presentation(pptx_stream)
        # Get filename from slide_data
        original_filename = next(iter(slide_data.values()))["filename"]

    else:
        raise ValueError("Either input_folder and output_folder or pptx_file_content must be provided.")

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_info = slide_data.get(slide_number, {})
        headline = slide_info.get("slide_headline", "")
        observations = slide_info.get("slide_observations", "")

        if not headline or headline == "HEADER SLIDE":
            logging.info(f"Slide {slide_number}: Skipped (Header or non-content slide)")
            continue

        # Update title placeholder with headline
        title_updated = False
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                shape.text = headline
                title_updated = True
                logging.info(f"Slide {slide_number}: Title updated with headline.")
                break

        if not title_updated:
            logging.warning(f"Slide {slide_number}: No title placeholder found for headline.")

        # Add observations to speaker notes
        if observations:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = observations
            logging.info(f"Slide {slide_number}: Observations added to speaker notes.")

    # Handle saving to disk
    if input_folder and output_folder:
        # Ensure the output directory exists
        os.makedirs(output_folder, exist_ok=True)

        # Save the modified presentation
        new_filename = original_filename.replace(".pptx", "_WITH_HEADLINES.pptx")
        new_pptx_path = os.path.join(output_folder, new_filename)

        presentation.save(new_pptx_path)
        logging.info(f"PowerPoint file saved with headlines and observations: {new_pptx_path}")

        return new_pptx_path

    # Handle returning bytes
    else:
        output_stream = BytesIO()
        presentation.save(output_stream)
        output_stream.seek(0)
        new_filename = original_filename.replace(".pptx", "_WITH_HEADLINES.pptx")

        logging.info(f"PowerPoint file prepared with headlines and observations as bytes: {new_filename}")

        return new_filename, output_stream.getvalue()
